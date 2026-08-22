"""Build a nat-os project overview deck.

Outputs:
  nat-os-project-overview.pptx
  ppt_assets/layers.png, memory.png, timeline.png

Run: platformio python make_presentation.py
"""

import os

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, Rectangle

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "ppt_assets")
os.makedirs(ASSETS, exist_ok=True)

NAVY = "#1F3864"
BLUE = "#2E74B5"
RED = "#C00000"
ORANGE = "#BF8F00"
GREEN = "#2E7D32"
GRAY = "#595959"
LIGHT = "#EDF2FA"

MONO = "Consolas"


# ---------------------------------------------------------------- figures --

def layer_stack():
    fig = plt.figure(figsize=(10.5, 6.2), dpi=200)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis("off")

    layers = [
        ("L4  Applications", "Bytecode VM, shell, launcher, NatScript assembler",
         "#A9D18E", "Written from scratch"),
        ("L3  Drivers & IPC", "Display, touch, flash, SD, ADC, I2C, audio, WiFi bridge",
         "#BDD7EE", "Written from scratch"),
        ("L2  Kernel", "Scheduler, heap, arenas, mutex, panic, clock",
         "#2F5597", "Written from scratch  (white text)"),
        ("L1  Bootloader", "Image load, flash MMU, partition table",
         "#F4B183", "Own 2,736 B loader (UM-NATOS-035)"),
        ("L0  ROM", "First-stage boot, crypto, SPI flash driver",
         "#D9D9D9", "Silicon — cannot replace"),
    ]
    x0, w = 0.12, 0.76
    y = 0.88
    h = 0.14
    gap = 0.025
    for title, desc, color, note in layers:
        ax.add_patch(FancyBboxPatch((x0, y - h), w, h,
                                    boxstyle="round,pad=0.006,rounding_size=0.012",
                                    facecolor=color, edgecolor=GRAY, lw=1.0))
        tc = "white" if color == "#2F5597" else NAVY
        ax.text(x0 + 0.02, y - h / 2 + 0.028, title, fontsize=11.5,
                fontweight="bold", color=tc, va="center")
        ax.text(x0 + 0.02, y - h / 2 - 0.018, desc, fontsize=8.8,
                color=tc if color == "#2F5597" else GRAY, va="center")
        ax.text(x0 + w + 0.02, y - h / 2, note, fontsize=8.5, color=GRAY,
                va="center", style="italic")
        y -= h + gap

    ax.text(0.5, 0.04,
            "Scope decision: start writing at L2. Only L0 is irreplaceable; "
            "L1 was replaced; clock config moved up to L2 (UM-NATOS-036).",
            fontsize=9, ha="center", color=NAVY, style="italic")

    p = os.path.join(ASSETS, "layers.png")
    fig.savefig(p, dpi=200, bbox_inches="tight")
    plt.close(fig)
    return p


def memory_overview():
    fig = plt.figure(figsize=(12.0, 5.6), dpi=200)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis("off")

    regions = [
        (0.72, "IRAM  0x40080000", "128 KB — kernel code, vectors, window.S bridges",
         "#2F5597", "All nat-os instructions lived here until UM-NATOS-037"),
        (0.52, "IROM  0x400D0000", "Flash-mapped executable — kernel .flash.text + vendor blob",
         "#5B9BD5", "WiFi build maps 606 KB pre-linked blob here"),
        (0.32, "DRAM  0x3FFB0000", "144 KB — .data, .bss, heap, task stacks, blob writable state",
         "#A9D18E", "No MMU paging — one address space for everything"),
        (0.12, "DROM  0x3F400020", "Flash-mapped read-only — .rodata since M11",
         "#F4B183", "64 KB page congruence; cache must stay on"),
    ]
    x0, w = 0.08, 0.55
    for y, title, body, color, note in regions:
        ax.add_patch(Rectangle((x0, y), w, 0.14, facecolor=color,
                               edgecolor=GRAY, lw=0.8, alpha=0.85))
        ax.text(x0 + 0.015, y + 0.105, title, fontsize=10.5, fontweight="bold",
                color="white" if color == "#2F5597" else NAVY)
        ax.text(x0 + 0.015, y + 0.055, body, fontsize=8.5, color=NAVY)
        ax.text(0.68, y + 0.07, note, fontsize=8.5, color=GRAY, va="center",
                style="italic")

    ax.text(0.5, 0.96, "ESP32 memory — no per-process virtual address spaces",
            fontsize=12, fontweight="bold", ha="center", color=NAVY)
    ax.text(0.5, 0.02,
            "Isolation is software: bytecode VM bounds-checks every load/store; "
            "per-app arenas; device permissions (UM-NATOS-032).",
            fontsize=9, ha="center", color=NAVY, style="italic")

    p = os.path.join(ASSETS, "memory.png")
    fig.savefig(p, dpi=200, bbox_inches="tight")
    plt.close(fig)
    return p


def milestone_timeline():
    fig = plt.figure(figsize=(12.4, 4.8), dpi=200)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis("off")

    milestones = [
        ("M0", "Boot & self-checks", "1,216 B image prints banner"),
        ("M1", "Timer & tick", "Level-3 interrupt, CCOMPARE1"),
        ("M2", "Preemptive tasks", "3,400+ switches, zero corruption"),
        ("M3", "Heap & arenas", "167 KB allocatable DRAM measured"),
        ("M4", "Bytecode VM", "35 opcodes, 6 fault classes contained"),
        ("M5", "Multi-app shell", "Escape attempt fails by design"),
    ]
    n = len(milestones)
    x_start, x_end = 0.06, 0.94
    ys = 0.55
    ax.plot([x_start, x_end], [ys, ys], color=NAVY, lw=2.5)

    for i, (tag, title, detail) in enumerate(milestones):
        x = x_start + i * (x_end - x_start) / (n - 1)
        ax.plot([x, x], [ys - 0.03, ys + 0.03], color=NAVY, lw=2)
        ax.add_patch(plt.Circle((x, ys), 0.018, facecolor=BLUE, edgecolor=NAVY, lw=1.2))
        ax.text(x, ys + 0.07, tag, fontsize=11, fontweight="bold", ha="center", color=NAVY)
        ax.text(x, ys + 0.13, title, fontsize=8.8, ha="center", color=NAVY)
        ax.text(x, ys - 0.08, detail, fontsize=7.5, ha="center", color=GRAY,
                wrap=True)

    # post-M5 bar
    ax.add_patch(FancyBboxPatch((0.06, 0.08), 0.88, 0.22,
                                boxstyle="round,pad=0.008,rounding_size=0.012",
                                facecolor=LIGHT, edgecolor=BLUE, lw=1.0))
    ax.text(0.5, 0.24,
            "After M5 (2026-08-14 → 2026-08-21): display DMA, touch, launcher, "
            "3D raycaster, persistence, SD, device model, custom bootloader, "
            "code-in-flash, vendor WiFi blob integration — 39 engineering reports",
            fontsize=9, ha="center", va="center", color=NAVY)
    ax.text(0.5, 0.12,
            "138 commits in three days, then a week of drivers, radio, and "
            "measurement-first debugging on windowed vendor code",
            fontsize=8.5, ha="center", va="center", color=GRAY, style="italic")

    p = os.path.join(ASSETS, "timeline.png")
    fig.savefig(p, dpi=200, bbox_inches="tight")
    plt.close(fig)
    return p


# ------------------------------------------------------------------ pptx ---

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_footer(slide, idx, total):
    tb = slide.shapes.add_textbox(Inches(0.35), Inches(7.08), Inches(9), Inches(0.35))
    tf = tb.text_frame
    r = tf.paragraphs[0].add_run()
    r.text = "nat-os project overview | Used Medias LLC — Embedded Systems Division"
    r.font.size = Pt(10); r.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
    tb2 = slide.shapes.add_textbox(Inches(12.35), Inches(7.08), Inches(0.7), Inches(0.35))
    tf2 = tb2.text_frame
    tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT
    r2 = tf2.paragraphs[0].add_run()
    r2.text = "%d / %d" % (idx, total)
    r2.font.size = Pt(10); r2.font.color.rgb = RGBColor(0x80, 0x80, 0x80)


def title_bar(slide, text, sub=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.92))
    shp.fill.solid(); shp.fill.fore_color.rgb = RGBColor(0x1F, 0x38, 0x64)
    shp.line.fill.background(); shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = Inches(0.4); tf.margin_top = Inches(0.08)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(26); r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    if sub:
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(12.5); r2.font.color.rgb = RGBColor(0xBD, 0xD7, 0xEE)


def bullets(slide, items, top=1.15, width=12.6, size=15):
    tb = slide.shapes.add_textbox(Inches(0.42), Inches(top), Inches(width),
                                  Inches(7.0 - top))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        lvl, text = item if isinstance(item, tuple) else (0, item)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(9)
        r = p.add_run()
        r.text = ("- " if lvl == 0 else "- ") + text
        r.font.size = Pt(size if lvl == 0 else size - 1.5)
        r.font.color.rgb = RGBColor(0x26, 0x26, 0x26)
    return tb


def content_slide(prs, title, sub=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(slide, title, sub)
    return slide


def table(slide, header, rows, widths, left=0.42, top=1.2, fsz=11.5,
          mono_cols=()):
    total_w = sum(widths)
    g = slide.shapes.add_table(len(rows) + 1, len(header), Inches(left),
                               Inches(top), Inches(total_w),
                               Inches(0.42 * (len(rows) + 1)))
    tbl = g.table
    for i, wd in enumerate(widths):
        tbl.columns[i].width = Inches(wd)
    hdr = tbl.rows[0]
    for j, htxt in enumerate(header):
        c = hdr.cells[j]
        c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0x1F, 0x38, 0x64)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]
        r = p.add_run(); r.text = htxt
        r.font.bold = True; r.font.size = Pt(fsz)
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    for i, row in enumerate(rows):
        tr = tbl.rows[i + 1]
        for j, val in enumerate(row):
            c = tr.cells[j]
            c.fill.solid()
            shade = 0xF2 if i % 2 == 0 else 0xFF
            c.fill.fore_color.rgb = RGBColor(shade, shade, 0xFF if shade == 0xFF else shade)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = Inches(0.06); c.margin_right = Inches(0.06)
            p = c.text_frame.paragraphs[0]
            r = p.add_run(); r.text = val
            r.font.size = Pt(fsz - 0.5)
            if j in mono_cols:
                r.font.name = MONO
                r.font.size = Pt(fsz - 1.5)
            if j == 0:
                r.font.bold = True
                r.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)
    return tbl


def place_picture(slide, img_path, top_in=1.05, max_h_in=5.85):
    pic = slide.shapes.add_picture(img_path, Inches(0), Inches(top_in))
    ratio = min(Inches(max_h_in) / pic.height, SLIDE_W / pic.width)
    pic.height = int(pic.height * ratio)
    pic.width = int(pic.width * ratio)
    pic.left = int((SLIDE_W - pic.width) / 2)
    return pic


def build():
    layers_png = layer_stack()
    mem_png = memory_overview()
    time_png = milestone_timeline()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    total = 16

    # 1 — title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x1F, 0x38, 0x64)
    bg.line.fill.background(); bg.shadow.inherit = False
    tb = s.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(11.8), Inches(3.2))
    tf = tb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = "nat-os"
    r.font.size = Pt(52); r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p = tf.add_paragraph(); p.space_before = Pt(8)
    r = p.add_run()
    r.text = "An operating system written from scratch for the ESP32"
    r.font.size = Pt(24); r.font.color.rgb = RGBColor(0xBD, 0xD7, 0xEE)
    p2 = tf.add_paragraph(); p2.space_before = Pt(28)
    r = p2.add_run()
    r.text = ("No ESP-IDF · No Arduino · No FreeRTOS · No C library\n"
              "ESP32-2432S028R (Cheap Yellow Display) · 39 engineering reports · 138 commits")
    r.font.size = Pt(15); r.font.color.rgb = RGBColor(0xBF, 0xBF, 0xBF)
    p3 = tf.add_paragraph(); p3.space_before = Pt(16)
    r = p3.add_run()
    r.text = "Project analysis — generated from source tree and docs/"
    r.font.size = Pt(14); r.font.color.rgb = RGBColor(0xE6, 0x91, 0x38)
    add_footer(s, 1, total)

    # 2 — executive summary
    s = content_slide(prs, "Executive summary")
    bullets(s, [
        "nat-os is a freestanding kernel for the ESP32 that owns scheduling, "
        "memory, drivers, and application execution. Every instruction from the "
        "image entry point onward is project code except a replaceable bootloader "
        "and partition table.",
        "Two decisions shape everything: compile the kernel with -mabi=call0 "
        "(no register windows — context switch is a conventional save/restore), "
        "and run applications in a bounds-checked bytecode VM (hardware has no "
        "MMU paging, so software isolation is the only option).",
        "Milestones M0–M5 completed in one day (2026-08-14): boot, tick, "
        "preemptive scheduler, heap/arenas, bytecode interpreter, multi-app shell. "
        "A week of driver work followed: display, touch, 3D raycaster, persistence, "
        "SD, audio, device model, custom bootloader, flash-executable code, and "
        "integration of Espressif's closed WiFi/PHY blob.",
        "The project is unusually well documented: 39 numbered engineering reports "
        "(UM-NATOS-001 through 039), a 31-chapter book synthesis, and standing rules "
        "extracted from every embarrassing defect.",
        "Current frontier: vendor WiFi stack runs from flash under nat-os's own "
        "bootloader; init-time windowed/call0 bridging bugs are being closed one "
        "measurement at a time. Nothing has been on air yet.",
    ], size=15.5)

    # 3 — hardware
    s = content_slide(prs, "Target hardware",
                      "ESP32-2432S028R — the only board every measurement was taken on")
    table(s,
          ["Property", "Value", "Design consequence"],
          [
              ["SoC", "ESP32-D0WD, dual Xtensa LX6 @ up to 240 MHz",
               "Real SMP possible; deferred — single-core kernel for now"],
              ["Internal SRAM", "520 KB total, no PSRAM on stock CYD",
               "Tight once graphics + WiFi blob writable state are reserved"],
              ["Flash", "4 MB, execute-in-place via cache/MMU",
               ".rodata in flash since M11; .text in flash since UM-NATOS-037"],
              ["MMU", "Flash-cache translation only",
               "No per-process address spaces — isolation is VM + arenas"],
              ["Display", "ILI9341 240×320 + XPT2046 resistive touch",
               "Board-specific code confined to board_cyd.h, display.c, touch.c"],
              ["Image size", "~37 KB kernel (no WiFi) / larger with blob",
               "~145 KB DRAM left for heap, stacks, and app arenas"],
          ],
          widths=[1.5, 4.2, 6.9],
          top=1.25, fsz=12)

    # 4 — architectural decisions
    s = content_slide(prs, "Two decisions that constrain everything")
    bullets(s, [
        "call0 ABI — Xtensa's windowed ABI (call8/call12) makes context switch "
        "require spilling live register windows — where from-scratch Xtensa kernels "
        "usually stall. call0 removes windows entirely. Cost: ROM routines (windowed) "
        "cannot be called directly — irrelevant until WiFi, when hand-written "
        "bridges in kernel/window.S become the integration layer.",
        "Bytecode VM for applications — The ESP32 MMU does not implement page "
        "tables. Any native code can write any address. Running apps in an "
        "interpreter whose loads/stores are bounds-checked recovers protection in "
        "software and makes preemption at instruction boundaries trivial. An app "
        "deliberately written to escape its arena is in the test suite — it cannot.",
        "L2 start — Writing L1 (clock trees, SPI calibration) is silicon bring-up, "
        "not OS design. Borrowing L1 yielded a running CPU quickly; L1 was later "
        "replaced entirely (UM-NATOS-035). Clock config moved to kernel/clock.c "
        "after the half-speed board defect (UM-NATOS-036).",
    ], size=15)

    # 5 — layer model
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "Layer model and project scope")
    place_picture(s, layers_png, top_in=1.05, max_h_in=5.9)
    add_footer(s, 5, total)

    # 6 — memory
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "Memory map and the isolation problem")
    place_picture(s, mem_png, top_in=1.05, max_h_in=5.9)
    add_footer(s, 6, total)

    # 7 — scheduler
    s = content_slide(prs, "Kernel: scheduler and tasks")
    bullets(s, [
        "Preemptive, three priority levels with ageing — no ready task waits more "
        "than ~600 ms regardless of what sits above it.",
        "Task states: READY, RUNNING, BLOCKED (mutex), SLEEPING (WAITI in idle). "
        "Up to 12 native task slots; context switch saves/restores a fixed frame "
        "in task.c / vectors.S.",
        "Tick: CCOMPARE1 level-3 interrupt. Critical standing rules emerged here: "
        "clear PS.EXCM before calling C from handlers (zero-overhead loops break "
        "otherwise); yield must never defer the tick deadline; shared register "
        "shadows need one writer or synchronized bookkeeping.",
        "Stack guards checked on every switch; worst task uses 444 B of 2,048 "
        "(78% margin measured). Panic dumps EXCCAUSE/EPC to serial, flash, and panel.",
        "WiFi build adds windowed vendor code via phy_stack_call, rom_callN, "
        "win_spill_all — the call0 kernel must bridge into -mabi=windowed blob code "
        "safely (kernel/window.S, vendor/windowed/wifi_osi_stubs.c).",
    ], size=14.5)

    # 8 — VM
    s = content_slide(prs, "Applications: bytecode VM and containment")
    bullets(s, [
        "Register-based VM: 35 opcodes, 14 syscalls, NatScript assembler in tools/. "
        "Programs compiled to .vasm, assembled to headers in kernel/generated/.",
        "Every memory access bounds-checked against the app's arena. Six fault "
        "classes contained (bad opcode, OOB load/store, stack overflow, etc.). "
        "3.3M instructions across 450 preemptions with zero accounting drift.",
        "Syscalls: display (FILL/TEXT/DIMS/BLIT confined to viewport), touch "
        "(confined), messaging (copied through kernel mailbox — never shared memory), "
        "device model (light, speaker, persistence, I2C, keypad, SD — with per-app "
        "permissions, UM-NATOS-032).",
        "Launcher: 3×3 icon grid, hybrid cursor, double-tap to run. On-screen shell "
        "reaches the real serial shell via UART tee — one command set, not two.",
        "Sample apps: pong, paint, rogue, raycast demo, note pad, spin, blit test.",
    ], size=14.5)

    # 9 — drivers table
    s = content_slide(prs, "Drivers and peripherals — verified state")
    table(s,
          ["Subsystem", "Hardware", "Status"],
          [
              ["Display", "ILI9341 via SPI2 + DMA", "43 ms full-screen fill; DMA bit-30 bug fixed (UM-NATOS-030)"],
              ["Touch", "XPT2046 + PENIRQ", "Calibrated on-device; X was inverted 3 months"],
              ["Audio", "LEDC PWM on GPIO26", "Tones work; 440 Hz inaudible, 3 kHz clear"],
              ["ADC", "SAR ADC1, 8 ch", "Light sensor on GPIO34 confirmed (+265 counts)"],
              ["I2C", "Bit-bang GPIO22/27", "Bus electrically verified; no device attached"],
              ["Flash store", "SPI flash @ 2 MB", "Checksummed record survives power cycle"],
              ["microSD", "SPI mode", "FAT16 header read; per-stage error codes"],
              ["Interrupts", "Matrix routing", "Peripheral IRQ routable; PENIRQ never fired from finger"],
              ["WiFi/PHY", "Espressif blob in flash", "PHY init OK; driver init in progress; not on air"],
              ["3D view", "Software raycaster", "16 fps grid dungeon; optional framebuffer"],
          ],
          widths=[1.6, 3.4, 7.6],
          top=1.2, fsz=11.5)

    # 10 — timeline
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "Development timeline",
              "138 commits · 2026-08-14 through 2026-08-21")
    place_picture(s, time_png, top_in=1.15, max_h_in=5.75)
    add_footer(s, 10, total)

    # 11 — build
    s = content_slide(prs, "Build pipeline")
    bullets(s, [
        "Toolchain: xtensa-esp32-elf-gcc from PlatformIO packages only — not "
        "PlatformIO as a build system. build.ps1 compiles, links, packages natos.bin.",
        "Key flags: -mabi=call0, -ffreestanding, -fno-builtin, -Os. WiFi build "
        "adds -DBOARD_WIFI_OVERRIDE=1 and links vendor blob + window bridges.",
        "Flash layout: stage-2 boot @ 0x1000, partition table @ 0x8000, kernel "
        "@ 0x10000. Custom bootloader is 2,736 bytes (boot/); vendor boot available "
        "via -VendorBootloader for A/B recovery.",
        "Board targets: cyd (default, verified) and lora32 (pin map not yet verified).",
        "Host tools: NatScript assembler (tools/vasm.py), serial command sender "
        "(tools/serial/send_cmd.py), bytecode program sources (*.vasm).",
    ], size=14.5)

    # 12 — documentation
    s = content_slide(prs, "Documentation culture",
                      "what makes this project unusual")
    bullets(s, [
        "39 engineering reports (UM-NATOS-001–039): each states what was decided, "
        "what was measured, and what remains unverified. Every report ends with "
        "what it does NOT establish.",
        "Measured vs assumed is always separated. Several standing rules exist "
        "because plausible-looking numbers were believed while an instrument "
        "declared its own reading invalid.",
        "docs/book/: 31-chapter synthesis with appendices (ISA reference, register "
        "map, timeline, measurements). Reports remain primary; book is continuity.",
        "Defects are recorded honestly: tick deadline racing into the future, "
        "touch axis inverted behind a calibration that could only give one answer, "
        "DMA_OUTLINK_START defined as bit 30 (RESTART not START), idle task that "
        "silently failed to be created.",
        "AGENTS.md codifies debugging discipline: reproduce first, one variable at "
        "a time, preserve known-good git state, treat linker maps and register dumps "
        "as primary evidence.",
    ], size=14)

    # 13 — standing rules
    s = content_slide(prs, "Standing rules (sample)",
                      "each rule constrains future code — extracted from a specific defect")
    table(s,
          ["Rule", "Origin", "One-line statement"],
          [
              ["EXCM", "M2 / UM-NATOS-009", "Clear PS.EXCM before calling C from any IRQ handler"],
              ["Tick yield", "Display freeze / UM-NATOS-016", "Yield must only move CCOMPARE1 earlier, never later"],
              ["Shadow registers", "183 ms stall / UM-NATOS-008", "One writer, or every writer maintains the shadow"],
              ["Instruments", "Touch cal / UM-NATOS-017", "When an instrument says its reading is invalid, believe it"],
              ["Experiments", "Flash debug / UM-NATOS-018", "A negative result is informative only if the experiment ran"],
              ["Existence ≠ proof", "Shell RX / UM-NATOS-019", "Trigger the mechanism on purpose, or treat as untested"],
              ["Lock contention", "Mutex timing / UM-NATOS-014", "Cost is number of blocking events, not hold time"],
              ["Fix scope", "Launcher chrome / UM-NATOS-021", "Correct diagnosis does not license arbitrary blast radius"],
          ],
          widths=[1.3, 3.2, 8.1],
          top=1.2, fsz=11)

    # 14 — recent work
    s = content_slide(prs, "Recent frontier work (reports 027–039)")
    bullets(s, [
        "UM-NATOS-028–030: WiFi receive path, seven instruments caught lying, "
        "display DMA dead since boot — root cause one wrong bit (RESTART vs START).",
        "UM-NATOS-031–032: Device model — peripherals as table entries with "
        "per-app permissions; ~70 lines, explicitly not security.",
        "UM-NATOS-034: Second receiver negative result — nothing demodulable on air; "
        "fault is RF/PHY path, not frame construction.",
        "UM-NATOS-035–037: Custom bootloader, half-speed clock defect, kernel code "
        "moved to flash (free IRAM 19 KB → 45 KB).",
        "UM-NATOS-038: 606 KB vendor blob mapped in three windows (code, rodata, "
        "writable DRAM); 118-entry OSI adapter table; struct layout off by 4 bytes "
        "found by asking real ESP-IDF to print offsets.",
        "UM-NATOS-039: Phantom WINDOWSTART bits root-caused — context switch wrote "
        "windowstart through rotated register view; fixed by instruction reorder. "
        "Exposed next failure: first multi-frame windowed task suspension.",
    ], size=13.5)

    # 15 — open problems
    s = content_slide(prs, "Current status and open problems")
    table(s,
          ["Area", "State", "Next / open"],
          [
              ["Core kernel M0–M5", "Complete, hardware-verified", "Stable baseline"],
              ["Display / touch / UI", "Working on CYD", "PENIRQ from finger unobserved"],
              ["WiFi transmit", "Not on air", "PHY init OK; driver init fault under investigation"],
              ["Windowed bridging", "Partially closed", "Multi-frame park machinery for blob task"],
              ["lora32 board", "Pin map unverified", "Second receiver experiment used CYD only"],
              ["Priority inheritance", "Unwired", "task_boost() exists; nothing calls it"],
              ["DRAM budget", "Unremeasured", "12 tasks + 3 drivers since last M3 figure"],
              ["JTAG probe", "Ordered, not in hand", "Would change debug methodology"],
          ],
          widths=[2.2, 3.8, 6.6],
          top=1.2, fsz=11.5)

    # 16 — source tree
    s = content_slide(prs, "Source tree at a glance")
    bullets(s, [
        "kernel/ — ~125 files: start.S, vectors.S, window.S (entry/IRQ/window bridges); "
        "task.c, timer.c, heap.c, arena.c, vm.c, app.c, ipc.c; display.c, touch.c, "
        "raycast.c, flash.c, sd.c, audio.c, i2c.c, adc.c; blobcall.c, phyinit.c, "
        "wifi_osi_impl.c; shell.c, term.c, desktop.c (launcher); linker.ld.",
        "boot/ — custom second-stage loader (replaces 17 KB ESP-IDF boot).",
        "vendor/ — partition table, optional ESP-IDF boot recovery, windowed WiFi blob "
        "and OSI stubs; net80211/ probe tooling.",
        "tools/ — NatScript assembler, demo/program .vasm sources, serial helpers.",
        "docs/ — 39 reports, book/, debug/ investigation records, next_moves/ roadmap.",
        "Key shell commands: ps, run, kill, mem, stacks, 3d, fb, sd, wifiinit (WiFi "
        "build), wintorture, hang/fault/smash (failure paths exercised on purpose).",
    ], size=13.5)

    # footers on content slides
    for i, slide in enumerate(prs.slides, start=1):
        if i == 1:
            continue
        has = any(shape.has_text_frame and "project overview |" in shape.text_frame.text
                  for shape in slide.shapes)
        if not has:
            add_footer(slide, i, total)

    out = os.path.join(HERE, "nat-os-project-overview.pptx")
    prs.save(out)
    return out


if __name__ == "__main__":
    out = build()
    print("WROTE:", out)
    print("SIZE:", os.path.getsize(out), "bytes")
